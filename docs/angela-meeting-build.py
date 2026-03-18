"""
Angela Supply Director Meeting — Presentation Builder
Run: python3 docs/angela-meeting-build.py
Output: docs/angela-supply-meeting-2026-03-18.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Brand Colors ────────────────────────────────────────────────────────────
NAVY       = RGBColor(0x00, 0x2B, 0x5E)   # dark navy — titles / BG
BLUE       = RGBColor(0x00, 0x5F, 0xB5)   # medium blue — accents
LIGHT_BLUE = RGBColor(0xE6, 0xF0, 0xFF)   # very light blue — subtle BG
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED        = RGBColor(0xD7, 0x26, 0x3D)   # alert / negative
AMBER      = RGBColor(0xF5, 0xA6, 0x23)   # warning
GREEN      = RGBColor(0x27, 0xAE, 0x60)   # positive
DARK_GREY  = RGBColor(0x2D, 0x2D, 0x2D)
MID_GREY   = RGBColor(0x6B, 0x6B, 0x6B)
LIGHT_GREY = RGBColor(0xF4, 0xF6, 0xF9)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

# ─────────────────────────────────────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = util.Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=DARK_GREY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def add_paragraph(tf, text, font_size=14, bold=False, color=DARK_GREY,
                  align=PP_ALIGN.LEFT, space_before=6, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    if text:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return p


def slide_header(slide, title, subtitle=None):
    """Standard navy top bar with title."""
    add_rect(slide, 0, 0, 13.33, 1.1, NAVY)
    add_text(slide, title, 0.4, 0.15, 10, 0.75,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.4, 0.75, 10, 0.35,
                 font_size=13, bold=False, color=RGBColor(0xAA, 0xC8, 0xFF),
                 align=PP_ALIGN.LEFT)
    # thin bottom accent line
    add_rect(slide, 0, 1.1, 13.33, 0.04, BLUE)


def kpi_box(slide, l, t, w, h, value, label, value_color=NAVY, bg=LIGHT_GREY,
            sub_label=None):
    add_rect(slide, l, t, w, h, bg, line_color=RGBColor(0xCC, 0xD8, 0xEE), line_width=0.5)
    # value
    add_text(slide, value, l + 0.1, t + 0.15, w - 0.2, h * 0.45,
             font_size=32, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    # label
    add_text(slide, label, l + 0.1, t + h * 0.5, w - 0.2, h * 0.3,
             font_size=11, bold=False, color=MID_GREY, align=PP_ALIGN.CENTER)
    if sub_label:
        add_text(slide, sub_label, l + 0.1, t + h * 0.78, w - 0.2, h * 0.2,
                 font_size=9, italic=True, color=MID_GREY, align=PP_ALIGN.CENTER)


def bar_chart_manual(slide, l, t, chart_w, chart_h, data, max_val,
                     bar_color=BLUE, label_color=DARK_GREY, show_values=True):
    """
    data = list of (label, value) tuples
    Draws horizontal bars manually.
    """
    n = len(data)
    row_h = chart_h / n
    bar_area_w = chart_w * 0.60
    label_w    = chart_w * 0.35

    for i, (label, val) in enumerate(data):
        y = t + i * row_h
        # Role label
        add_text(slide, label,
                 l, y + row_h * 0.18,
                 label_w - 0.05, row_h * 0.65,
                 font_size=9, color=label_color, align=PP_ALIGN.RIGHT)
        # Bar background
        add_rect(slide, l + label_w, y + row_h * 0.2,
                 bar_area_w, row_h * 0.58,
                 RGBColor(0xE0, 0xE8, 0xF5))
        # Bar fill
        bar_w = bar_area_w * (val / max_val)
        add_rect(slide, l + label_w, y + row_h * 0.2,
                 bar_w, row_h * 0.58, bar_color)
        # Value label
        if show_values:
            add_text(slide, str(val),
                     l + label_w + bar_w + 0.05,
                     y + row_h * 0.2,
                     0.3, row_h * 0.58,
                     font_size=9, bold=True, color=label_color)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — TITLE
# ─────────────────────────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, NAVY)
add_rect(s1, 0, 5.8, 13.33, 1.7, BLUE)

# Decorative accent block
add_rect(s1, 0, 0, 0.25, 7.5, BLUE)

add_text(s1, "Supply Performance Review", 0.6, 1.4, 12, 1.0,
         font_size=44, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "AI Interview Design   |   Funnel Impact   |   Optimization Opportunities",
         0.6, 2.55, 12, 0.6,
         font_size=18, color=RGBColor(0xAA, 0xC8, 0xFF), align=PP_ALIGN.LEFT)

add_rect(s1, 0.6, 3.3, 4.0, 0.06, BLUE)

add_text(s1, "Angela - Supply Director", 0.6, 3.55, 8, 0.4,
         font_size=14, color=RGBColor(0xCC, 0xDD, 0xFF), align=PP_ALIGN.LEFT)
add_text(s1, "March 18, 2026", 0.6, 3.95, 8, 0.35,
         font_size=13, italic=True, color=RGBColor(0x88, 0xAA, 0xDD), align=PP_ALIGN.LEFT)

add_text(s1, "Confidential - Internal Use Only", 1.0, 6.1, 11.33, 0.5,
         font_size=11, italic=True, color=RGBColor(0xCC, 0xDD, 0xFF),
         align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — FUNNEL BASELINE (Context setter)
# ─────────────────────────────────────────────────────────────────────────────
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.5, WHITE)
slide_header(s2, "The Worker Funnel - 2025 Baseline",
             "52,987 accounts created > 8,146 first shifts completed (15.4% end-to-end)")

# 5 KPI boxes
kpi_data = [
    ("52,987", "Accounts Created",    NAVY,  None),
    ("56.0%",  "Account > Verified",  RED,   "Biggest drop-off"),
    ("36.0%",  "Reach Platform Verified", AMBER, "Cumulative"),
    ("18.2%",  "Book a 1st Shift",    AMBER, "Cumulative"),
    ("15.4%",  "Complete 1st Shift",  GREEN, "End-to-end goal"),
]
kpi_x = [0.3, 2.9, 5.5, 8.1, 10.7]
for i, (val, lbl, col, sub) in enumerate(kpi_data):
    kpi_box(s2, kpi_x[i], 1.3, 2.25, 1.6, val, lbl, value_color=col,
            sub_label=sub)

# Arrow connectors (simple text arrows)
for x in [2.58, 5.18, 7.78, 10.38]:
    add_text(s2, ">", x, 1.75, 0.3, 0.5, font_size=20, bold=True, color=MID_GREY)

# Funnel stage breakdown table
add_rect(s2, 0.3, 3.15, 12.73, 0.35, NAVY)
headers = ["Funnel Stage", "Count", "Stage CR%", "Cumulative", "People Lost"]
col_x   = [0.4, 4.5, 6.2, 8.1, 10.1]
col_w   = [4.0, 1.6, 1.8, 1.8, 2.9]
for i, h in enumerate(headers):
    add_text(s2, h, col_x[i], 3.18, col_w[i], 0.28,
             font_size=10, bold=True, color=WHITE)

rows = [
    ("Worker Accounts Created",    "52,987", "-",     "100%",  "-"),
    ("1st Role Verified",          "29,693", "56.0%", "56.0%", "23,294  << AI Interview drop"),
    ("1st OB Task Completed",      "27,614", "93.0%", "52.1%", "2,079"),
    ("Platform Verified",          "19,091", "69.1%", "36.0%", "8,523"),
    ("1st Shift Booked",           " 9,628", "63.2%", "18.2%", "9,463"),
    ("1st Shift Completed",        " 8,146", "84.6%", "15.4%", "1,482"),
]
row_colors = [LIGHT_GREY, RGBColor(0xFF, 0xEB, 0xED), LIGHT_GREY, RGBColor(0xFF, 0xF3, 0xE0),
              LIGHT_GREY, RGBColor(0xE8, 0xF8, 0xEE)]
for r, (row, bg) in enumerate(zip(rows, row_colors)):
    y = 3.5 + r * 0.55
    add_rect(s2, 0.3, y, 12.73, 0.52, bg)
    txt_colors = [DARK_GREY, DARK_GREY, DARK_GREY, DARK_GREY, RED if r == 1 else MID_GREY]
    for i, cell in enumerate(row):
        add_text(s2, cell, col_x[i], y + 0.08, col_w[i], 0.38,
                 font_size=9.5, color=txt_colors[i],
                 bold=(r == 1 and i == 4))

add_text(s2, "Source: OB Funnel Custom Viewer (44).xlsx - Full Year 2025, 147 clients",
         0.3, 7.1, 12, 0.3, font_size=8, italic=True, color=MID_GREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — AI INTERVIEW: THE HIDDEN FRICTION
# ─────────────────────────────────────────────────────────────────────────────
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.5, WHITE)
slide_header(s3, "AI Interview: Hidden Friction at the Top of the Funnel",
             "Question count varies 2x across roles - directly linked to Account > Verified drop-off")

# Left panel — insight text
add_rect(s3, 0.3, 1.3, 5.8, 5.9, LIGHT_GREY)
add_rect(s3, 0.3, 1.3, 0.08, 5.9, RED)  # red left border

txb = s3.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(5.3), Inches(5.5))
txb.word_wrap = True
tf = txb.text_frame
tf.word_wrap = True

p0 = tf.paragraphs[0]
p0.alignment = PP_ALIGN.LEFT
r0 = p0.add_run()
r0.text = "The Problem"
r0.font.size = Pt(15)
r0.font.bold = True
r0.font.color.rgb = RED

add_paragraph(tf, "44% of workers who create an account never complete role verification - "
              "the AI interview is the primary barrier.", 12, color=DARK_GREY, space_before=8)
add_paragraph(tf, "", space_before=4)

p_head2 = tf.add_paragraph()
r_h2 = p_head2.add_run()
r_h2.text = "Why Question Volume Matters"
r_h2.font.size = Pt(14)
r_h2.font.bold = True
r_h2.font.color.rgb = NAVY

facts = [
    "Shortest interview: 16 questions (Warehouse Clerk)",
    "Longest interview: 32 questions (Housekeeper) - 2x longer",
    "Hospitality avg = 22 Qs vs. Industrial avg = 19 Qs",
    "Hospitality clients convert 40-50% lower than Industrial",
    "Bartender (30 Qs) and Server (29 Qs) = near-max friction",
    "Most hospitality questions are screening-only (0 scored) - no differentiation value",
]
for fact in facts:
    add_paragraph(tf, "  " + fact, 11, color=DARK_GREY, space_before=5)

add_paragraph(tf, "", space_before=4)

p_head3 = tf.add_paragraph()
r_h3 = p_head3.add_run()
r_h3.text = "Conversion Gap - Hospitality vs. Industrial"
r_h3.font.size = Pt(12)
r_h3.font.bold = True
r_h3.font.color.rgb = NAVY

comp_rows = [
    ("Merritt Hospitality",  "7.5%"),
    ("Legends Hospitality",  "8.8%"),
    ("LowCountry Catering",  "11.9%"),
    ("-- vs. ----------------", "------"),
    ("Tennant Solutions",    "20.3%"),
    ("CORT (Industrial)",    "19.0%"),
    ("AFC Industries",       "14.3%"),
]
for label, val in comp_rows:
    p = tf.add_paragraph()
    p.space_before = Pt(3)
    r = p.add_run()
    r.text = f"  {label:<28} {val}"
    r.font.size = Pt(10)
    is_hosp = "Hospitality" in label or "Catering" in label
    is_sep   = "--" in label
    r.font.color.rgb = RED if is_hosp else (MID_GREY if is_sep else GREEN)
    r.font.bold = not is_sep

# Right panel — bar chart
add_text(s3, "Interview Length by Role (# Questions)", 6.5, 1.3, 6.5, 0.4,
         font_size=13, bold=True, color=NAVY)

industrial_data = [
    ("Warehouse Clerk",       16),
    ("Warehouse Operative",   18),
    ("Loader Crew",           17),
    ("Picker Packer",         18),
    ("Assembler",             18),
    ("Material Handler",      19),
    ("Machine Operator",      20),
    ("Repair Technician",     20),
    ("Industrial General Lab",23),
    ("Forklift Driver",       21),
]
hospitality_data = [
    ("Event Staff",           17),
    ("Dishwasher",            17),
    ("Busser",                18),
    ("Hosp General Labor",    18),
    ("Concession Stand",      19),
    ("Bar Back",              20),
    ("Host",                  21),
    ("Barista",               22),
    ("Prep Cook",             22),
    ("Line Cook",             23),
    ("Hosp Team Lead",        24),
    ("Server",                29),
    ("Bartender",             30),
    ("Housekeeper",           32),
]

# Legend
add_rect(s3, 6.5, 1.75, 0.18, 0.18, BLUE)
add_text(s3, "Industrial", 6.72, 1.73, 2, 0.22, font_size=9, color=DARK_GREY)
add_rect(s3, 8.5, 1.75, 0.18, 0.18, RED)
add_text(s3, "Hospitality", 8.72, 1.73, 2, 0.22, font_size=9, color=DARK_GREY)

# Avg line reference text
add_text(s3, "| Avg 19", 10.5, 1.73, 1.2, 0.22, font_size=8, italic=True, color=BLUE)
add_text(s3, "| Avg 22", 11.5, 1.73, 1.2, 0.22, font_size=8, italic=True, color=RED)

# Industrial bars
bar_chart_manual(s3, 6.3, 2.0, 6.7, 2.55, industrial_data, 35,
                 bar_color=BLUE, label_color=DARK_GREY)
# Hospitality bars
bar_chart_manual(s3, 6.3, 4.65, 6.7, 2.55, hospitality_data, 35,
                 bar_color=RED, label_color=DARK_GREY)

add_text(s3, "Source: INDUSTRIAL TACO UPDATES.xlsx + HOSPITALITY TACO UPDATES.xlsx - TACO system - 2026-03-17",
         0.3, 7.1, 12.7, 0.3, font_size=8, italic=True, color=MID_GREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — RECOMMENDATIONS
# ─────────────────────────────────────────────────────────────────────────────
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.5, WHITE)
slide_header(s4, "Recommendations - Reducing Interview Drop-Off",
             "Priority actions to recover the 23,294 workers lost at Account > Verified in 2025")

recs = [
    (
        "01", "Audit & Slim High-Question Roles", BLUE,
        [
            "Housekeeper (32), Bartender (30), Server (29) are 50-100% longer than the minimum.",
            "Propose a cap: 20 questions per role for screening-only interviews (0 scored).",
            "Target: reduce hospitality average from 22 to 18 questions.",
            "Expected impact: faster time-to-complete, lower abandonment rate.",
        ]
    ),
    (
        "02", "Score More, Screen Less", AMBER,
        [
            "13 of 25 roles have zero scored questions - every question is a blocker, not a filter.",
            "Converting screening Qs into scored ones keeps length but adds quality signal.",
            "Priority: Server, Bartender, Housekeeper - high volume, high drop-off.",
            "Scored interviews also give recruiters ranking data to prioritize callbacks.",
        ]
    ),
    (
        "03", "Set Expectation on the Landing Page", NAVY,
        [
            "Workers who know the interview exists upfront don't abandon mid-flow.",
            "Add '10-minute video interview' language to role-specific landing pages.",
            "Specifically: highest drop-off roles (Housekeeper, Bartender, Server).",
            "Benchmark: this tactic alone can recover 5-10pp at Account > Verified.",
        ]
    ),
    (
        "04", "Track Interview Completion Rate by Role", GREEN,
        [
            "Currently the funnel measures Account > Verified as a single step.",
            "Instrument: 'Interview Started' and 'Interview Completed' as distinct events.",
            "This isolates whether drop-off is pre-interview (LP/UX) or mid-interview (length).",
            "Data needed before any interview redesign can be validated.",
        ]
    ),
]

box_positions = [(0.3, 1.25), (6.85, 1.25), (0.3, 4.25), (6.85, 4.25)]

for (num, title, color, bullets), (bx, by) in zip(recs, box_positions):
    bw, bh = 6.2, 2.85
    add_rect(s4, bx, by, bw, bh, LIGHT_GREY,
             line_color=RGBColor(0xCC, 0xD8, 0xEE), line_width=0.5)
    add_rect(s4, bx, by, 0.5, bh, color)  # colored left bar

    # Number badge
    add_rect(s4, bx + 0.55, by + 0.15, 0.38, 0.38, color)
    add_text(s4, num, bx + 0.55, by + 0.15, 0.38, 0.38,
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(s4, title, bx + 1.0, by + 0.15, bw - 1.1, 0.38,
             font_size=13, bold=True, color=NAVY)

    txb2 = s4.shapes.add_textbox(
        Inches(bx + 0.65), Inches(by + 0.65), Inches(bw - 0.8), Inches(bh - 0.75)
    )
    txb2.word_wrap = True
    tf2 = txb2.text_frame
    tf2.word_wrap = True
    for j, b in enumerate(bullets):
        p = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = "  " + b
        r.font.size = Pt(10)
        r.font.color.rgb = DARK_GREY


# ─────────────────────────────────────────────────────────────────────────────
# Helper: Table builder
# ─────────────────────────────────────────────────────────────────────────────

def add_table_slide(slide, headers, rows, col_widths, top_y,
                    header_bg=NAVY, header_fg=WHITE,
                    row_colors=None, cell_colors=None):
    """Draw a table manually using rectangles and text boxes."""
    left_margin = 0.3
    row_h = 0.42
    hdr_h = 0.38
    total_w = sum(col_widths)

    # Header row
    add_rect(slide, left_margin, top_y, total_w, hdr_h, header_bg)
    x = left_margin
    for i, hdr in enumerate(headers):
        add_text(slide, hdr, x + 0.08, top_y + 0.04, col_widths[i] - 0.16, hdr_h - 0.08,
                 font_size=9, bold=True, color=header_fg)
        x += col_widths[i]

    # Data rows
    for r, row_data in enumerate(rows):
        y = top_y + hdr_h + r * row_h
        bg = (row_colors[r] if row_colors else
              (LIGHT_GREY if r % 2 == 0 else WHITE))
        add_rect(slide, left_margin, y, total_w, row_h, bg,
                 line_color=RGBColor(0xDD, 0xDD, 0xDD), line_width=0.3)
        x = left_margin
        for i, cell in enumerate(row_data):
            fg = DARK_GREY
            bld = False
            if cell_colors and (r, i) in cell_colors:
                fg = cell_colors[(r, i)]
                bld = True
            add_text(slide, str(cell), x + 0.08, y + 0.06, col_widths[i] - 0.16, row_h - 0.12,
                     font_size=9, color=fg, bold=bld)
            x += col_widths[i]


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — CAMPAIGN PERFORMANCE KPIs
# ─────────────────────────────────────────────────────────────────────────────
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.5, WHITE)
slide_header(s5, "Indeed Campaign Performance — March 2026 MTD",
             "Campaign spend efficiency across all active client markets")

# KPI boxes
kpi5 = [
    ("$39,683", "Total Spend (Matched)", NAVY, None),
    ("2,583", "Total RSVPs", NAVY, None),
    ("$15.36", "Avg Cost/RSVP", BLUE, "Benchmark"),
    ("16.1%", "Apply > RSVP", AMBER, "Conversion rate"),
    ("$0.61", "Avg CPC", GREEN, None),
]
kpi5_x = [0.3, 2.9, 5.5, 8.1, 10.7]
for i, (val, lbl, col, sub) in enumerate(kpi5):
    kpi_box(s5, kpi5_x[i], 1.3, 2.25, 1.5, val, lbl, value_color=col, sub_label=sub)

# Top Performers table
add_text(s5, "Top Performers (lowest Cost/RSVP)", 0.3, 3.0, 6.2, 0.4,
         font_size=13, bold=True, color=GREEN)

top_headers = ["Client", "Location", "Cost/RSVP", "Apply>RSVP", "RSVPs", "Fill %"]
top_widths = [1.6, 1.5, 1.0, 1.0, 0.8, 0.7]
top_rows = [
    ("Legends Hospitality", "Dallas, TX", "$3.40", "54.3%", "114/275", "41%"),
    ("Stord", "Las Vegas, NV", "$3.63", "55.1%", "209/960", "22%"),
    ("CTDI", "Grove City, OH", "$4.48", "69.6%", "80/400", "20%"),
    ("CORT", "Houston, TX", "$4.88", "49.1%", "82/300", "27%"),
    ("Stord", "Atlanta, GA", "$5.57", "33.7%", "122/400", "30%"),
    ("CORT", "Chicago, IL", "$7.57", "29.5%", "167/400", "42%"),
    ("OnTrac", "Orlando, FL", "$7.84", "38.2%", "39/60", "65%"),
]
# Color the Cost/RSVP column green
top_colors = {}
for r in range(len(top_rows)):
    top_colors[(r, 2)] = GREEN

add_table_slide(s5, top_headers, top_rows, top_widths, 3.4,
                cell_colors=top_colors)

# Worst Performers table
add_text(s5, "Worst Performers (highest Cost/RSVP)", 6.95, 3.0, 6.2, 0.4,
         font_size=13, bold=True, color=RED)

worst_headers = ["Client", "Location", "Cost/RSVP", "App>RSVP", "Issue"]
worst_widths = [1.1, 1.5, 0.9, 0.8, 2.1]
worst_rows = [
    ("OnTrac", "S. Brunswick, NJ", "$1,103", "0.3%", "$4.4K spend, 4 RSVPs"),
    ("Culinaire", "Dallas, TX", "$429", "0.7%", "$3K spend, 7 RSVPs"),
    ("CTDI", "Haslet, TX", "$93", "2.4%", "$4.3K spend, 46 RSVPs"),
    ("Tennant", "Cincinnati, OH", "$84", "4.0%", "$336 spend, 4 RSVPs"),
    ("CORT", "Nashville, TN", "$79", "3.4%", "$1.1K spend, 14 RSVPs"),
    ("Cont. Battery", "Reno, NV", "$79", "3.3%", "$867 spend, 11 RSVPs"),
    ("CORT", "Las Vegas, NV", "$68", "3.1%", "$1.6K spend, 23 RSVPs"),
]
worst_colors = {}
for r in range(len(worst_rows)):
    worst_colors[(r, 2)] = RED

add_table_slide(s5, worst_headers, worst_rows, worst_widths, 3.4,
                header_bg=RGBColor(0x8B, 0x00, 0x00),
                cell_colors=worst_colors)

add_text(s5, "Source: Indeed Analytics (JobsCampaigns) + FHS Requisitions — March 1-17, 2026",
         0.3, 7.1, 12, 0.3, font_size=8, italic=True, color=MID_GREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — WASTED SPEND + KEY FINDINGS
# ─────────────────────────────────────────────────────────────────────────────
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, 13.33, 7.5, WHITE)
slide_header(s6, "Spend Leakage & Key Findings",
             "~$3.2K going to locations with no open requisition + critical funnel issues")

# Wasted spend table (left)
add_text(s6, "Spend with No Open Requisition", 0.3, 1.3, 6.2, 0.4,
         font_size=13, bold=True, color=RED)
add_text(s6, "These campaigns are spending money with no matching job in FHS",
         0.3, 1.65, 6.2, 0.3, font_size=10, italic=True, color=MID_GREY)

waste_headers = ["Client", "Location", "Spend", "Apply Starts", "Note"]
waste_widths = [1.2, 1.4, 0.8, 0.9, 2.1]
waste_rows = [
    ("Ingram", "La Vergne, TN", "$1,056", "527", "No open req found"),
    ("Johnstone", "Lancaster, TX", "$500", "216", "No open req found"),
    ("CORT", "Rhino, NC", "$490", "218", "No open req found"),
    ("Evergreen", "Austin/Dallas", "$481", "216", "Brand campaigns"),
    ("BTX", "Austin, TX", "$231", "89", "No open req found"),
    ("Hyatt", "Atlanta, GA", "$227", "101", "No open req found"),
]
waste_colors = {}
for r in range(len(waste_rows)):
    waste_colors[(r, 2)] = RED

add_table_slide(s6, waste_headers, waste_rows, waste_widths, 2.0,
                cell_colors=waste_colors)

add_rect(s6, 0.3, 4.7, 6.3, 0.45, RGBColor(0xFF, 0xEB, 0xED))
add_text(s6, "Total wasted: ~$2,985 on campaigns with no matching open requisition",
         0.5, 4.75, 6.0, 0.35, font_size=10, bold=True, color=RED)

# Key Findings (right panel)
add_rect(s6, 7.0, 1.3, 6.0, 5.6, LIGHT_GREY,
         line_color=RGBColor(0xCC, 0xD8, 0xEE), line_width=0.5)
add_rect(s6, 7.0, 1.3, 0.08, 5.6, NAVY)

txb6 = s6.shapes.add_textbox(Inches(7.25), Inches(1.45), Inches(5.6), Inches(5.3))
txb6.word_wrap = True
tf6 = txb6.text_frame
tf6.word_wrap = True

p_title = tf6.paragraphs[0]
r_t = p_title.add_run()
r_t.text = "Key Findings"
r_t.font.size = Pt(16)
r_t.font.bold = True
r_t.font.color.rgb = NAVY

findings = [
    ("1", "OnTrac South Brunswick burning $4.4K for 4 RSVPs",
     "Apply>RSVP rate is 0.3% — 1,309 people start applying but almost none RSVP. "
     "Something is broken in the flow for this location.", RED),
    ("2", "CTDI Haslet absorbs 8.4% of total spend",
     "$4.3K with only 2.4% Apply>RSVP — worst ratio among high-spend locations.", AMBER),
    ("3", "Efficiency winners under $6/RSVP",
     "Legends ($3.40), Stord LV ($3.63), CTDI Grove City ($4.48), CORT Houston ($4.88) — "
     "all with 30-70% Apply>RSVP rates.", GREEN),
    ("4", "~$3K going to locations with no open req",
     "Ingram, Johnstone, BTX, Hyatt, Evergreen — potential wasted spend "
     "or campaigns not yet linked to requisitions.", RED),
]

for num, title, detail, color in findings:
    add_paragraph(tf6, "", space_before=10)
    p_n = tf6.add_paragraph()
    p_n.space_before = Pt(2)
    r_num = p_n.add_run()
    r_num.text = f"  {num}. {title}"
    r_num.font.size = Pt(11)
    r_num.font.bold = True
    r_num.font.color.rgb = color

    p_d = tf6.add_paragraph()
    p_d.space_before = Pt(2)
    r_detail = p_d.add_run()
    r_detail.text = f"     {detail}"
    r_detail.font.size = Pt(9)
    r_detail.font.color.rgb = DARK_GREY

add_text(s6, "Source: Indeed Analytics + FHS Requisitions — March 1-17, 2026",
         0.3, 7.1, 12, 0.3, font_size=8, italic=True, color=MID_GREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — LAS VEGAS HIRING EVENTS RESULTS
# ─────────────────────────────────────────────────────────────────────────────
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, 13.33, 7.5, WHITE)
slide_header(s7, "Las Vegas Hiring Events — March 11 & 16, 2026",
             "2 events held | 187 invited | 89 arrived (48%) | 91 workers cleared for shifts")

# --- KPI Row ---
kpi7 = [
    ("187", "Total Invited", NAVY, "Across 2 events"),
    ("89", "Total Arrived", BLUE, "48% show rate"),
    ("85", "DT Passed", GREEN, "96% of arrivals"),
    ("91", "Workers Cleared", GREEN, "Incl. 6 DT-only"),
    ("52%", "No-Show Rate", RED, "98 did not arrive"),
]
kpi7_x = [0.3, 2.9, 5.5, 8.1, 10.7]
for i, (val, lbl, col, sub) in enumerate(kpi7):
    kpi_box(s7, kpi7_x[i], 1.3, 2.25, 1.4, val, lbl, value_color=col, sub_label=sub)

# --- Left: Event Comparison Table ---
add_text(s7, "Event Comparison", 0.3, 2.9, 6.2, 0.35,
         font_size=13, bold=True, color=NAVY)

evt_headers = ["Metric", "Mar 11 Event", "Mar 16 Event", "Combined"]
evt_widths = [1.8, 1.4, 1.4, 1.4]
evt_rows = [
    ("Invited", "104", "83", "187"),
    ("Arrived", "43 (41%)", "46 (55%)", "89 (48%)"),
    ("DT Passed", "40 (93%)", "45 (98%)", "85 (96%)"),
    ("DT Failed", "3", "1", "4"),
    ("ACP Completed", "—", "44 (96%)", "—"),
    ("Skilled (noted)", "36", "—", "—"),
    ("No-Shows", "61 (59%)", "37 (45%)", "98 (52%)"),
]
evt_colors = {
    (1, 2): GREEN,  # Mar 16 better show rate
    (2, 2): GREEN,  # Mar 16 better DT rate
    (6, 0): RED,    # No-shows label
}
add_table_slide(s7, evt_headers, evt_rows, evt_widths, 3.25, cell_colors=evt_colors)

# --- Client Preference (Mar 16) ---
add_text(s7, "Client Preference (Mar 16)", 0.3, 6.6, 3, 0.3,
         font_size=10, bold=True, color=NAVY)
add_text(s7, "Stord: 34  |  CORT: 6  |  Both: 2", 0.3, 6.85, 3, 0.25,
         font_size=9, color=DARK_GREY)

# --- Right: Recruiter Feedback Panel ---
add_rect(s7, 6.55, 2.9, 6.45, 4.3, LIGHT_GREY,
         line_color=RGBColor(0xCC, 0xD8, 0xEE), line_width=0.5)
add_rect(s7, 6.55, 2.9, 0.08, 4.3, AMBER)

txb7 = s7.shapes.add_textbox(Inches(6.8), Inches(3.0), Inches(6.0), Inches(4.0))
txb7.word_wrap = True
tf7 = txb7.text_frame
tf7.word_wrap = True

p_ft = tf7.paragraphs[0]
r_ft = p_ft.add_run()
r_ft.text = "Recruiter Feedback (Marcus)"
r_ft.font.size = Pt(14)
r_ft.font.bold = True
r_ft.font.color.rgb = NAVY

feedback_items = [
    ("AI Interview Blocker", RED,
     "Flexers had scheduled interviews blocking on-the-spot AI interviews. "
     "~12 workers waited while workaround was found (cancel via email → reopen in app). "
     "Recruiting team had to conduct in-person interviews to prevent further delays."),
    ("Interview Expectation Gap", AMBER,
     "Almost every flexer expected an interview at the event, defeating the purpose "
     "of AI interviews (complete anytime at convenience). Recommend: require interview "
     "completion BEFORE event or advise clearly in invitation."),
    ("Venue & Logistics", BLUE,
     "Address change likely caused lower attendance. No Indeed Flex signage present. "
     "Meeting rooms used for privacy, but some flexers preferred car/outside. "
     "Need dedicated interview space for future events."),
    ("Process Improvement", GREEN,
     "Google Sheet check-in worked well — recommend pre-building for all events. "
     "Only ~2 people came specifically for DT. Document verification team needs "
     "company phone number (not personal devices)."),
]

for title, color, detail in feedback_items:
    add_paragraph(tf7, "", space_before=6)
    p_i = tf7.add_paragraph()
    p_i.space_before = Pt(2)
    r_i = p_i.add_run()
    r_i.text = f"  {title}"
    r_i.font.size = Pt(10)
    r_i.font.bold = True
    r_i.font.color.rgb = color

    p_d = tf7.add_paragraph()
    p_d.space_before = Pt(1)
    r_d = p_d.add_run()
    r_d.text = f"  {detail}"
    r_d.font.size = Pt(8)
    r_d.font.color.rgb = DARK_GREY

add_text(s7, "Source: Check In Hiring Event.xlsx + Recruiter feedback (Marcus) — March 2026",
         0.3, 7.15, 12, 0.25, font_size=8, italic=True, color=MID_GREY)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out = "/Users/claudio.santos/RM-Team-Ai/docs/angela-supply-meeting-2026-03-18.pptx"
prs.save(out)
print(f"Saved: {out}")
