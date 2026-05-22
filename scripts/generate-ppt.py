#!/usr/bin/env python3
"""
Paid Ads Performance Review — PowerPoint Generator
Period: April 22 – May 22, 2026
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Brand Colors ──────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x0F, 0x2B, 0x5B)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GREEN       = RGBColor(0x00, 0xA8, 0x78)
RED         = RGBColor(0xE6, 0x39, 0x46)
AMBER       = RGBColor(0xF4, 0xA2, 0x61)
LIGHT_GRAY  = RGBColor(0xF5, 0xF7, 0xFA)
MID_GRAY    = RGBColor(0xD0, 0xD5, 0xDD)
DARK_GRAY   = RGBColor(0x34, 0x3A, 0x40)
LIGHT_BLUE  = RGBColor(0xE8, 0xEF, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=12, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color or DARK_GRAY
    return txb


def header_bar(slide, title):
    add_rect(slide, 0, 0, 13.33, 0.7, fill=DARK_BLUE)
    add_text(slide, title, 0.3, 0.08, 12.5, 0.55,
             font_size=18, bold=True, color=WHITE)


def slide_bg(slide):
    add_rect(slide, 0, 0, 13.33, 7.5, fill=LIGHT_GRAY)


def callout_box(slide, text, x, y, w, h, bg=GREEN, icon=""):
    add_rect(slide, x, y, w, h, fill=bg)
    add_text(slide, icon + text, x + 0.12, y + 0.05, w - 0.2, h - 0.1,
             font_size=10, bold=False, color=WHITE, wrap=True)


def section_label(slide, text, x, y, w=4):
    add_rect(slide, x, y, w, 0.28, fill=DARK_BLUE)
    add_text(slide, text, x + 0.08, y + 0.03, w - 0.1, 0.22,
             font_size=9, bold=True, color=WHITE)


def kpi_card(slide, label, value, x, y, w=2.0, h=0.9, accent=DARK_BLUE):
    add_rect(slide, x, y, w, h, fill=WHITE, line_color=MID_GRAY, line_width=0.5)
    add_rect(slide, x, y, w, 0.07, fill=accent)
    add_text(slide, label, x + 0.1, y + 0.1, w - 0.15, 0.28,
             font_size=9, bold=False, color=DARK_GRAY)
    add_text(slide, value, x + 0.1, y + 0.35, w - 0.15, 0.45,
             font_size=16, bold=True, color=DARK_BLUE)


def table_header_row(slide, cols, x, y, col_widths, row_h=0.28):
    cx = x
    for i, col in enumerate(cols):
        add_rect(slide, cx, y, col_widths[i], row_h, fill=DARK_BLUE)
        add_text(slide, col, cx + 0.05, y + 0.04, col_widths[i] - 0.08, row_h - 0.06,
                 font_size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        cx += col_widths[i]


def table_data_row(slide, cols, x, y, col_widths, row_h=0.26, even=True,
                   cell_colors=None):
    cx = x
    bg = LIGHT_GRAY if even else WHITE
    for i, col in enumerate(cols):
        cell_bg = cell_colors[i] if cell_colors and cell_colors[i] else bg
        add_rect(slide, cx, y, col_widths[i], row_h, fill=cell_bg,
                 line_color=MID_GRAY, line_width=0.3)
        add_text(slide, str(col), cx + 0.05, y + 0.04, col_widths[i] - 0.08,
                 row_h - 0.06, font_size=8, bold=False, color=DARK_GRAY,
                 align=PP_ALIGN.CENTER)
        cx += col_widths[i]


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK_BLUE)

# Decorative accent bar
add_rect(slide, 0, 5.8, 13.33, 0.08, fill=GREEN)
add_rect(slide, 0, 5.9, 13.33, 0.05, fill=AMBER)

add_text(slide, "Paid Ads Performance Review",
         1.0, 1.8, 11.0, 1.2, font_size=40, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(slide, "Google  ·  Indeed  ·  Reddit",
         1.0, 3.1, 11.0, 0.7, font_size=22, bold=False, color=GREEN,
         align=PP_ALIGN.CENTER)
add_text(slide, "April 22 – May 22, 2026",
         1.0, 3.9, 11.0, 0.5, font_size=16, bold=False, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(slide, "Prepared by: Recruitment Marketing Team",
         1.0, 4.5, 11.0, 0.4, font_size=12, bold=False, color=MID_GRAY,
         align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — EXECUTIVE SUMMARY
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Executive Summary — Total Spend & Channel Mix")

# 3 KPI cards
kpi_card(slide, "Total Spend (Period)", "$189,229", 0.3, 0.85, 2.8, 1.0, DARK_BLUE)
kpi_card(slide, "Prior Period Total", "~$241,150", 3.3, 0.85, 2.8, 1.0, DARK_GRAY)
kpi_card(slide, "Month-over-Month Change", "–21% ▼", 6.3, 0.85, 2.8, 1.0, RED)

# Donut-style channel share visual (3 colored boxes)
share_data = [
    ("Indeed Ads", "$111,257", "58.8%", DARK_BLUE),
    ("Google Ads", "$69,823",  "36.9%", GREEN),
    ("Reddit Ads", "$8,149",   "4.3%",  AMBER),
]
bx = 9.3
for label, spend, pct, color in share_data:
    add_rect(slide, bx, 0.85, 3.7 * (float(pct.rstrip('%')) / 100), 0.45, fill=color)
    bx += 3.7 * (float(pct.rstrip('%')) / 100)

bx = 9.3
for label, spend, pct, color in share_data:
    bar_w = 3.7 * (float(pct.rstrip('%')) / 100)
    add_text(slide, pct, bx + 0.04, 0.88, bar_w, 0.38,
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    bx += bar_w

add_text(slide, "Channel Breakdown", 9.3, 1.35, 3.7, 0.25, font_size=9,
         bold=True, color=DARK_BLUE)
for i, (label, spend, pct, color) in enumerate(share_data):
    add_rect(slide, 9.3, 1.62 + i * 0.22, 0.16, 0.14, fill=color)
    add_text(slide, f"{label}  {pct}  ({spend})", 9.5, 1.60 + i * 0.22, 3.6, 0.2,
             font_size=8.5, bold=False, color=DARK_GRAY)

# Table
headers  = ["Channel", "Spend", "Share", "Prior Period", "MoM Change"]
col_w    = [2.2, 1.6, 1.1, 1.9, 1.6]
table_x  = 0.3
table_y  = 2.05
table_header_row(slide, headers, table_x, table_y, col_w)

rows = [
    ("Indeed Ads",  "$111,257", "58.8%", "~$137,000", "–19%"),
    ("Google Ads",  "$69,823",  "36.9%", "$93,834",   "–26%"),
    ("Reddit Ads",  "$8,149",   "4.3%",  "$10,318",   "–21%"),
    ("TOTAL",       "$189,229", "100%",  "~$241,150",  "–21%"),
]
for i, row in enumerate(rows):
    bold_bg = LIGHT_BLUE if row[0] == "TOTAL" else None
    cc = [bold_bg] * 5 if bold_bg else None
    table_data_row(slide, row, table_x, table_y + 0.28 + i * 0.26,
                   col_w, even=(i % 2 == 0), cell_colors=cc)

callout_box(slide,
    "All three channels declining. Total budget down $52K month-over-month.",
    0.3, 3.45, 12.7, 0.52, bg=RED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — GOOGLE ADS OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Google Ads — Performance Overview  |  Apr 22 – May 22, 2026")

kpis = [
    ("Total Spend",       "$69,823"),
    ("Impressions",       "2,224,715"),
    ("Clicks",            "77,110"),
    ("Conversions",       "39,265"),
    ("Avg CPC",           "$0.91"),
    ("CPA",               "$1.78"),
    ("CTR",               "3.47%"),
    ("Active Campaigns",  "112"),
]
kw = 3.0
for i, (label, val) in enumerate(kpis):
    row, col = divmod(i, 4)
    kpi_card(slide, label, val, 0.3 + col * (kw + 0.08),
             0.85 + row * 1.05, kw, 0.92)

headers = ["Metric", "April (Full)", "May 1–21", "Change"]
col_w   = [2.5, 2.5, 2.5, 2.5]
table_y = 3.1
table_header_row(slide, headers, 0.3, table_y, col_w)

rows_g = [
    ("Spend",       "$93,834",   "$43,163",   "–54%"),
    ("Clicks",      "135,851",   "47,469",    "–65%"),
    ("Impressions", "3,264,575", "1,404,473", "–57%"),
    ("Conversions", "55,116",    "23,633",    "–57%"),
    ("CPA",         "$1.70",     "$1.83",     "+7.3%"),
    ("CPC",         "$0.69",     "$0.91",     "+31.6% ⚠"),
]
for i, row in enumerate(rows_g):
    cc = [None, None, None, RED if "+" in row[3] else None]
    table_data_row(slide, row, 0.3, table_y + 0.28 + i * 0.26, col_w,
                   even=(i % 2 == 0), cell_colors=cc)

callout_box(slide,
    "CPC rising +31.6% as volume falls — efficiency degrading in May.",
    0.3, 6.85, 12.7, 0.45, bg=RED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — GOOGLE ADS: CAMPAIGN TYPE BREAKDOWN
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Google Ads — Campaign Type Breakdown")

section_label(slide, "Performance by Campaign Type", 0.3, 0.82, 6.0)

headers = ["Type", "Spend", "Share", "Clicks", "Conversions", "CPA", "Signal"]
col_w   = [1.5, 1.4, 1.0, 1.4, 1.6, 1.2, 2.2]
table_y = 1.15
table_header_row(slide, headers, 0.3, table_y, col_w)

type_rows = [
    ("App",     "$17,673", "25.3%", "27,803", "19,539", "$0.90",  "✅ Most Efficient"),
    ("Search",  "$25,017", "35.8%", "17,039", "10,327", "$2.42",  "Moderate"),
    ("PMax",    "$18,179", "26.0%", "19,067", "6,805",  "$2.67",  "Moderate"),
    ("Display", "$4,158",  "6.0%",  "8,631",  "1,434",  "$2.90",  "Weak"),
    ("Other",   "$4,795",  "6.9%",  "4,570",  "1,160",  "$4.13",  "⚠ Inefficient"),
]
cpa_colors_map = {"$0.90": GREEN, "$2.42": None, "$2.67": None, "$2.90": None, "$4.13": AMBER}

for i, row in enumerate(type_rows):
    cpa_c = cpa_colors_map.get(row[5])
    cc = [None, None, None, None, None, cpa_c, None]
    table_data_row(slide, row, 0.3, table_y + 0.28 + i * 0.26, col_w,
                   even=(i % 2 == 0), cell_colors=cc)

# Bar chart (spend)
chart_x, chart_y, chart_w, chart_h = 0.3, 2.85, 8.0, 2.8
add_rect(slide, chart_x, chart_y, chart_w, chart_h, fill=WHITE,
         line_color=MID_GRAY, line_width=0.5)
add_text(slide, "Spend Distribution by Campaign Type", chart_x + 0.1,
         chart_y + 0.05, chart_w - 0.2, 0.28, font_size=10, bold=True, color=DARK_BLUE)

bar_data = [("App", 17673, GREEN), ("Search", 25017, DARK_BLUE),
            ("PMax", 18179, AMBER), ("Display", 4158, MID_GRAY), ("Other", 4795, RED)]
max_spend = max(d[1] for d in bar_data)
bar_area_w = chart_w - 1.5
bar_area_x = chart_x + 1.2
bar_y_start = chart_y + 0.5
bar_slot_h  = 0.36
for i, (lbl, spend, color) in enumerate(bar_data):
    by = bar_y_start + i * (bar_slot_h + 0.08)
    bw = bar_area_w * (spend / max_spend)
    add_text(slide, lbl, chart_x + 0.08, by + 0.04, 1.1, 0.28,
             font_size=8.5, bold=False, color=DARK_GRAY)
    add_rect(slide, bar_area_x, by, bw, bar_slot_h, fill=color)
    add_text(slide, f"${spend:,}", bar_area_x + bw + 0.05, by + 0.06,
             0.9, 0.24, font_size=8, bold=False, color=DARK_GRAY)

callout_box(slide,
    "App Campaigns deliver 2.7× better CPA than Search at $0.90 — highest ROI format.",
    0.3, 6.3, 12.7, 0.52, bg=GREEN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — GOOGLE ADS: PERFORMANCE BY MARKET
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Google Ads — Performance by Market")

headers = ["Market", "Spend", "Conversions", "CPA", "Rating"]
col_w   = [3.5, 1.8, 1.8, 1.4, 4.5]
table_y = 0.85
table_header_row(slide, headers, 0.3, table_y, col_w)

market_rows = [
    ("Chicago, IL",         "$9,307",  "4,685", "$1.99", "✅ Strong"),
    ("Dallas, TX",          "$8,334",  "2,730", "$3.05", "⚠ High CPA"),
    ("South Brunswick, NJ", "$4,694",  "2,899", "$1.62", "✅ Strong"),
    ("Paulsboro, NJ",       "$4,467",  "2,517", "$1.77", "✅ Strong"),
    ("Nashville, TN",       "$3,691",  "1,893", "$1.95", "✅ Solid"),
    ("Las Vegas, NV",       "$3,005",  "2,203", "$1.36", "✅ Efficient"),
    ("Austin, TX",          "$2,955",  "1,228", "$2.41", "Moderate"),
    ("Orlando, FL",         "$2,498",  "1,885", "$1.33", "✅ Efficient"),
    ("Reno, NV",            "$2,045",  "429",   "$4.77", "🔴 Struggling"),
    ("Atlanta, GA",         "$1,561",  "1,894", "$0.82", "✅ Best CPA"),
]
cpa_color_map = {
    "$1.99": GREEN, "$3.05": AMBER, "$1.62": GREEN, "$1.77": GREEN,
    "$1.95": GREEN, "$1.36": GREEN, "$2.41": None,  "$1.33": GREEN,
    "$4.77": RED,   "$0.82": GREEN,
}
for i, row in enumerate(market_rows):
    cpa_c = cpa_color_map.get(row[3])
    cc = [None, None, None, cpa_c, None]
    table_data_row(slide, row, 0.3, table_y + 0.28 + i * 0.26, col_w,
                   even=(i % 2 == 0), cell_colors=cc)

callout_box(slide, "Atlanta at $0.82 CPA — best market in the account.",
            0.3, 6.2, 6.1, 0.45, bg=GREEN)
callout_box(slide, "Reno at $4.77 CPA — 2.7× above target. Review keywords and audience.",
            6.55, 6.2, 6.5, 0.45, bg=RED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GOOGLE ADS: WEEKLY SPEND TREND
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Google Ads — Weekly Spend Trend")

weeks = [
    ("Mar 30", 15483, 8208),
    ("Apr 6",  22150, 13390),
    ("Apr 13", 22601, 14074),
    ("Apr 20", 22395, 12639),
    ("Apr 27", 20377, 11440),
    ("May 4",  16770, 9180),
    ("May 11", 10756, 6291),
    ("May 18", 7073,  3819),
]

chart_x, chart_y = 0.5, 0.85
chart_w, chart_h = 12.3, 4.8
max_spend = 25000
bar_count  = len(weeks)
bar_slot_w = chart_w / bar_count
bar_margin = 0.12

add_rect(slide, chart_x, chart_y, chart_w, chart_h, fill=WHITE,
         line_color=MID_GRAY, line_width=0.5)

# Grid lines
for pct in [0.25, 0.5, 0.75, 1.0]:
    gy = chart_y + chart_h - chart_h * pct
    add_rect(slide, chart_x, gy, chart_w, 0.005, fill=MID_GRAY)
    add_text(slide, f"${int(max_spend * pct):,}",
             chart_x - 1.1, gy - 0.1, 1.0, 0.22,
             font_size=8, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

for i, (week, spend, conv) in enumerate(weeks):
    bx  = chart_x + i * bar_slot_w + bar_margin
    bw  = bar_slot_w - bar_margin * 2
    bh  = (spend / max_spend) * chart_h * 0.92
    by  = chart_y + chart_h - bh
    color = RED if i == len(weeks) - 1 else (GREEN if spend == max([w[1] for w in weeks]) else DARK_BLUE)
    add_rect(slide, bx, by, bw, bh, fill=color)
    add_text(slide, f"${spend // 1000}K", bx, by - 0.28, bw, 0.25,
             font_size=8.5, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
    add_text(slide, week, bx - 0.05, chart_y + chart_h + 0.05,
             bw + 0.1, 0.28, font_size=8, color=DARK_GRAY, align=PP_ALIGN.CENTER)

callout_box(slide,
    "Spend dropped 69% from peak (Apr 13: $22,601) to current week (May 18: $7,073). "
    "Daily rate: ~$1,000/day now vs $3,200/day at peak.",
    0.3, 6.2, 12.7, 0.55, bg=RED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — GOOGLE ADS: TOP PERFORMERS vs GAPS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Google Ads — Top Performers vs Struggling Campaigns")

# LEFT column
section_label(slide, "✅ Top 5 Most Efficient", 0.3, 0.82, 6.1)
headers_l = ["Market", "Type", "Spend", "Conv", "CPA"]
col_w_l   = [2.6, 1.1, 1.0, 0.9, 0.9]
table_header_row(slide, headers_l, 0.3, 1.13, col_w_l)
top5 = [
    ("Atlanta, GA",            "App", "$461",  "1,136", "$0.41"),
    ("Charlotte, NC",          "App", "$237",  "577",   "$0.41"),
    ("Atlanta, GA",            "Sch", "$94",   "209",   "$0.45"),
    ("Middleburg Heights, OH", "App", "$249",  "515",   "$0.48"),
    ("Flower Mound, TX",       "App", "$925",  "1,541", "$0.60"),
]
for i, row in enumerate(top5):
    cc = [None, None, None, None, GREEN]
    table_data_row(slide, row, 0.3, 1.41 + i * 0.26, col_w_l,
                   even=(i % 2 == 0), cell_colors=cc)

# RIGHT column
section_label(slide, "🔴 Top 5 Struggling Campaigns", 6.75, 0.82, 6.2)
headers_r = ["Market", "Type", "Spend", "Conv", "CPA"]
col_w_r   = [2.6, 1.1, 1.0, 0.9, 0.9]
table_header_row(slide, headers_r, 6.75, 1.13, col_w_r)
bot5 = [
    ("Nashville, TN",  "Search", "$907",  "78",  "$11.63"),
    ("Chicago, IL",    "Search", "$666",  "54",  "$12.36"),
    ("Hebron, KY",     "Search", "$603",  "44",  "$13.64"),
    ("Competitor",     "Search", "$593",  "0",   "—"),
    ("Dallas, TX",     "D.Gen",  "$604",  "74",  "$8.16"),
]
for i, row in enumerate(bot5):
    cc = [None, None, None, None, RED]
    table_data_row(slide, row, 6.75, 1.41 + i * 0.26, col_w_r,
                   even=(i % 2 == 0), cell_colors=cc)

# Divider
add_rect(slide, 6.6, 0.82, 0.04, 4.2, fill=DARK_BLUE)

callout_box(slide,
    "$2,769 spent on 4 campaigns at 5–8× target CPA. Recommend pause or restructure.",
    0.3, 6.2, 12.7, 0.52, bg=RED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — INDEED ADS OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Indeed Ads — Performance Overview  |  Apr 23 – May 22, 2026")

kpis_i = [
    ("Total Spend",          "$111,257"),
    ("Est. Clicks",          "~98,500"),
    ("Apply Starts",         "~22,300"),
    ("Avg CPC",              "~$1.13"),
    ("Cost Per Apply Start", "~$4.98"),
]
kw = 2.45
for i, (label, val) in enumerate(kpis_i):
    kpi_card(slide, label, val, 0.3 + i * (kw + 0.1), 0.85, kw, 0.95)

# Note box
add_rect(slide, 0.3, 1.95, 12.7, 0.45, fill=LIGHT_BLUE, line_color=DARK_BLUE, line_width=0.5)
add_text(slide,
    "Note: Apply Starts = primary conversion metric on Indeed. "
    "Platform-tracked Applies are a secondary signal with historically low completion rates (~2–4%).",
    0.45, 1.98, 12.3, 0.38, font_size=8.5, italic=True, color=DARK_BLUE)

headers = ["Metric", "April Daily Avg", "May Daily Rate", "Change"]
col_w   = [3.0, 2.8, 2.8, 3.6]
table_y = 2.55
table_header_row(slide, headers, 0.3, table_y, col_w)

rows_i = [
    ("Spend",        "$4,582/day", "$2,760/day", "–40%"),
    ("Clicks",       "6,539/day",  "4,393/day",  "–33%"),
    ("Apply Starts", "1,618/day",  "1,034/day",  "–36%"),
    ("Avg CPC",      "$0.70",      "$0.63",      "–10% ✅"),
]
for i, row in enumerate(rows_i):
    cc = [None, None, None, GREEN if "✅" in row[3] else None]
    table_data_row(slide, row, 0.3, table_y + 0.28 + i * 0.26, col_w,
                   even=(i % 2 == 0), cell_colors=cc)

callout_box(slide,
    "Daily spend pace down 40% MoM — but CPC improved 10%. Apply Start volume tracking volume decline.",
    0.3, 4.2, 12.7, 0.52, bg=AMBER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — INDEED ADS: BY STATE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Indeed Ads — Performance by State  |  May 2026")

headers = ["State", "Spend", "Clicks", "Apply Starts", "Flag"]
col_w   = [2.6, 1.8, 1.6, 1.8, 5.5]
table_y = 0.85
table_header_row(slide, headers, 0.3, table_y, col_w)

state_rows = [
    ("Texas",          "$22,070", "40,788", "7,685",  "✅ Largest market"),
    ("Ohio",           "$8,470",  "12,367", "3,087",  "✅ Solid"),
    ("New Jersey",     "$4,911",  "7,576",  "1,882",  "⚠ Review"),
    ("Georgia",        "$4,385",  "4,428",  "1,692",  "⚠ Review"),
    ("DC / Washington","$3,539",  "5,006",  "1,658",  "⚠ Review"),
    ("Tennessee",      "$2,826",  "4,213",  "952",    "⚠ Review"),
    ("Illinois",       "$2,047",  "3,915",  "859",    "⚠ Review"),
    ("Nevada",         "$1,506",  "1,651",  "694",    "⚠ Review"),
    ("Florida",        "$1,470",  "1,686",  "557",    "⚠ Review"),
]
for i, row in enumerate(state_rows):
    cc = [None, None, None, None, GREEN if "✅" in row[4] else AMBER]
    table_data_row(slide, row, 0.3, table_y + 0.28 + i * 0.26, col_w,
                   even=(i % 2 == 0), cell_colors=cc)

callout_box(slide,
    "7 states show strong click volume but low platform-tracked apply completions — "
    "possible attribution gap or supply-side mismatch.",
    0.3, 6.3, 12.7, 0.52, bg=AMBER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — INDEED ADS: TOP CLIENTS & STRUGGLING ACCOUNTS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Indeed Ads — Top Clients & Attribution Note  |  May 2026")

section_label(slide, "Top Clients by Spend", 0.3, 0.82, 5.8)
headers_l = ["Client", "Spend", "Clicks", "Apply Starts"]
col_w_l   = [2.4, 1.5, 1.4, 1.5]
table_header_row(slide, headers_l, 0.3, 1.12, col_w_l)

top_clients = [
    ("Ontrac",          "$6,454",  "10,184", "2,474"),
    ("CTDI",            "$4,624",  "6,716",  "1,502"),
    ("CORT (combined)", "$3,721",  "4,225",  "1,487"),
    ("Stord",           "$3,265",  "3,420",  "1,313"),
    ("Indeed Flex BAU", "$3,212",  "7,075",  "1,455"),
    ("Legends",         "$2,296",  "2,669",  "851"),
    ("HEI",             "$1,456",  "1,685",  "543"),
]
for i, row in enumerate(top_clients):
    table_data_row(slide, row, 0.3, 1.40 + i * 0.26, col_w_l, even=(i % 2 == 0))

section_label(slide, "⚠ Platform Apply Attribution — Known Issue", 7.0, 0.82, 6.1)

add_rect(slide, 7.0, 1.12, 6.1, 3.5, fill=WHITE, line_color=AMBER, line_width=1)
note_lines = [
    "Several high-spend accounts show 0 platform-tracked",
    "applies despite strong Apply Start counts:",
    "",
    "  • Ontrac:  $6,454 spend → 2,474 Apply Starts, 0 applies",
    "  • CTDI:    $4,624 spend → 1,502 Apply Starts, 0 applies",
    "  • Stord:   $3,265 spend → 1,313 Apply Starts, 0 applies",
    "  • Legends: $2,296 spend → 851 Apply Starts, 0 applies",
    "",
    "0 applies ≠ 0 hires.",
    "",
    "Indeed's apply-tracking has known attribution gaps for",
    "external ATS flows. Apply Starts remain the reliable",
    "intent signal and primary optimization metric.",
]
for j, line in enumerate(note_lines):
    add_text(slide, line, 7.1, 1.18 + j * 0.24, 5.9, 0.22,
             font_size=8.5, bold=("≠" in line or "0 applies" in line),
             color=DARK_GRAY if "•" not in line else DARK_BLUE)

callout_box(slide,
    "Apply Starts are the primary performance signal on Indeed. Platform-tracked Applies "
    "have known ATS attribution gaps — do not use as the sole KPI.",
    0.3, 6.3, 12.7, 0.52, bg=AMBER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — REDDIT ADS OVERVIEW
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Reddit Ads — Performance Overview  |  Apr 22 – May 21, 2026")

kpis_r = [
    ("Total Spend",     "$8,149"),
    ("Impressions",     "620,582"),
    ("Clicks",          "3,661"),
    ("Results",         "14,305"),
    ("CPC",             "$2.23"),
    ("Cost Per Result", "$0.57"),
    ("eCPM",            "$12.62"),
    ("Active Campaigns","10"),
]
kw = 3.0
for i, (label, val) in enumerate(kpis_r):
    row, col = divmod(i, 4)
    kpi_card(slide, label, val, 0.3 + col * (kw + 0.08),
             0.85 + row * 1.05, kw, 0.92)

headers = ["Metric", "Mar 22–Apr 21", "Apr 22–May 21", "Change"]
col_w   = [2.8, 2.8, 2.8, 3.0]
table_y = 3.05
table_header_row(slide, headers, 0.3, table_y, col_w)

rows_r = [
    ("Spend",       "$10,318", "$8,149",  "–21%"),
    ("Impressions", "756,470", "620,582", "–18%"),
    ("Clicks",      "4,290",   "3,661",   "–15%"),
    ("Results",     "8,438",   "14,305",  "+70% ✅"),
    ("CPC",         "$2.41",   "$2.23",   "–7% ✅"),
    ("CPR",         "$1.22",   "$0.57",   "–53% ✅"),
]
for i, row in enumerate(rows_r):
    cc = [None, None, None, GREEN if "✅" in row[3] else None]
    table_data_row(slide, row, 0.3, table_y + 0.28 + i * 0.26, col_w,
                   even=(i % 2 == 0), cell_colors=cc)

callout_box(slide,
    "Results up 70% while spend fell 21% — CPR cut in half. "
    "Driven by FIFA Legends Concession campaign ($0.02 CPR).",
    0.3, 6.3, 12.7, 0.52, bg=GREEN)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — REDDIT ADS: CAMPAIGN BREAKDOWN
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Reddit Ads — Campaign Breakdown")

headers = ["Campaign", "Spend", "Impr", "Clicks", "Results", "CPR", "Flag"]
col_w   = [3.8, 1.1, 1.3, 1.1, 1.2, 1.0, 4.4]
table_y = 0.85
table_header_row(slide, headers, 0.3, table_y, col_w)

reddit_rows = [
    ("PP – McCarran, NV (Industrial)",   "$2,070", "139,395", "710",   "0",      "—",     "🔴 Zero Results"),
    ("Culinaire – Dallas (Hospitality)", "$1,939", "159,266", "959",   "0",      "—",     "🔴 Zero Results"),
    ("Ontrac Evergreen – Lebanon, OH",   "$902",   "59,398",  "241",   "150",    "$6.01", "⚠ High CPR"),
    ("Dishwasher – Nashville, TN",       "$894",   "71,157",  "540",   "441",    "$2.03", "✅ Good"),
    ("Legends Line Cook – Dallas, TX",   "$891",   "69,706",  "495",   "303",    "$2.94", "✅ Good"),
    ("Hospitality Evergreen – Austin",   "$656",   "52,569",  "311",   "254",    "$2.58", "✅ Good"),
    ("FIFA Legends Concession – Dallas", "$291",   "30,190",  "112",   "12,944", "$0.02", "✅ Best CPR"),
    ("Levy Concession – Dallas, TX",     "$257",   "18,400",  "144",   "111",    "$2.32", "✅ Good"),
    ("CORT Loader – Orlando, FL",        "$206",   "16,331",  "126",   "102",    "$2.02", "✅ Good"),
]
cpr_colors = {
    "—":     RED,
    "$6.01": AMBER,
    "$2.03": GREEN,
    "$2.94": GREEN,
    "$2.58": GREEN,
    "$0.02": GREEN,
    "$2.32": GREEN,
    "$2.02": GREEN,
}
for i, row in enumerate(reddit_rows):
    cpr_c = cpr_colors.get(row[5])
    cc = [None, None, None, None, None, cpr_c, RED if "🔴" in row[6] else (AMBER if "⚠" in row[6] else GREEN if "✅" in row[6] else None)]
    table_data_row(slide, row, 0.3, table_y + 0.28 + i * 0.26, col_w,
                   even=(i % 2 == 0), cell_colors=cc)

callout_box(slide,
    "$4,009 spent on McCarran + Culinaire with ZERO results. "
    "Combined waste across both periods: $11,260.",
    0.3, 6.2, 12.7, 0.52, bg=RED)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — GAPS & KEY ISSUES
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Gaps & Key Issues — All Channels")

# 3 columns
cols = [
    ("Google Ads", [
        "Nashville Search (Hospitality): $907, CPA $11.63 — 6.5× target",
        "Chicago Search (Hospitality): $666, CPA $12.36 — 6.9× target",
        "Hebron/Stord Forklift Search: $603, CPA $13.64 — 7.6× target",
        "Competitor Search: $593, 0 conversions",
        "Reno market: $2,045, CPA $4.77 — 2.7× above target",
        "CPC rising +31.6% MoM as volume shrinks",
        "Spend down 69% from Apr 13 peak — no floor visible",
    ]),
    ("Indeed Ads", [
        "Dallas: $11,335 spend, only 65 platform applies (0.7%)",
        "Washington DC: $3,539 spend, 0 platform applies",
        "Atlanta: $4,385 spend, 0 platform applies",
        "7 states show near-zero apply completion rate",
        "Apply completion rate ~2% overall — below industry norms",
        "Daily spend pace down 40% MoM, no recovery signal",
    ]),
    ("Reddit Ads", [
        "McCarran NV (Industrial): $6,233 total across periods — 0 results ever",
        "Culinaire Dallas: $5,027 total — 0 results ever",
        "Channel nearly paused: 3 campaigns at $25/week as of May 15–22",
        "No Industrial campaigns with positive ROI this period",
    ]),
]

for ci, (ch, issues) in enumerate(cols):
    cx = 0.25 + ci * 4.38
    add_rect(slide, cx, 0.82, 4.2, 0.32, fill=RED)
    add_text(slide, ch, cx + 0.1, 0.85, 4.0, 0.26,
             font_size=10, bold=True, color=WHITE)
    add_rect(slide, cx, 1.17, 4.2, 5.6, fill=WHITE,
             line_color=MID_GRAY, line_width=0.5)
    for ji, issue in enumerate(issues):
        add_text(slide, "• " + issue, cx + 0.1, 1.22 + ji * 0.7, 4.0, 0.65,
                 font_size=8.5, color=DARK_GRAY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — BUDGET ALLOCATION & MAY TREND
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Budget Allocation — April Actuals vs May Projected")

headers = ["Channel", "April Actual", "May Projected", "Change"]
col_w   = [2.5, 2.5, 2.5, 2.5]
table_y = 0.85
table_header_row(slide, headers, 0.3, table_y, col_w)

budget_rows = [
    ("Indeed",  "$137,451", "~$83,000",  "–40%"),
    ("Google",  "$93,834",  "~$60,000",  "–36%"),
    ("Reddit",  "$10,318",  "~$8,149",   "–21%"),
    ("TOTAL",   "$241,603", "~$151,149", "–37%"),
]
for i, row in enumerate(budget_rows):
    bold_bg = LIGHT_BLUE if row[0] == "TOTAL" else None
    cc = [bold_bg] * 4 if bold_bg else [None, None, None, RED]
    table_data_row(slide, row, 0.3, table_y + 0.28 + i * 0.26, col_w,
                   even=(i % 2 == 0), cell_colors=cc)

# Stacked bar comparison
section_label(slide, "April vs May Projected — Stacked Channel Spend", 0.3, 2.6, 7.0)
bar_groups = [
    ("April",     137451, 93834, 10318),
    ("May (proj)",83000,  60000, 8149),
]
max_total = 250000
bx = 1.5
for g, (label, indeed, google, reddit) in enumerate(bar_groups):
    bar_y = 3.0
    bar_w = 2.5
    for amount, color in [(indeed, DARK_BLUE), (google, GREEN), (reddit, AMBER)]:
        bh = (amount / max_total) * 3.2
        bar_y_pos = bar_y + (3.2 - (indeed + google + reddit) / max_total * 3.2)
        bar_y_pos = 6.2 - bh - sum(
            amt / max_total * 3.2
            for amt, _ in [(indeed, None), (google, None), (reddit, None)]
            if _ is None
        )
    # simplified stacked bars
    total = indeed + google + reddit
    by = 6.2
    for amt, col in [(reddit, AMBER), (google, GREEN), (indeed, DARK_BLUE)]:
        bh = (amt / max_total) * 3.2
        add_rect(slide, bx + g * 3.2, by - bh, bar_w, bh, fill=col)
        if bh > 0.2:
            add_text(slide, f"${amt//1000}K", bx + g * 3.2 + 0.05,
                     by - bh + 0.04, bar_w - 0.1, bh - 0.06,
                     font_size=8, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
        by -= bh
    add_text(slide, label, bx + g * 3.2, 6.25, bar_w, 0.28,
             font_size=9, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)

# Legend
add_rect(slide, 8.0, 3.0, 0.3, 0.2, fill=DARK_BLUE)
add_text(slide, "Indeed", 8.35, 3.0, 1.5, 0.2, font_size=8.5, color=DARK_GRAY)
add_rect(slide, 8.0, 3.3, 0.3, 0.2, fill=GREEN)
add_text(slide, "Google", 8.35, 3.3, 1.5, 0.2, font_size=8.5, color=DARK_GRAY)
add_rect(slide, 8.0, 3.6, 0.3, 0.2, fill=AMBER)
add_text(slide, "Reddit", 8.35, 3.6, 1.5, 0.2, font_size=8.5, color=DARK_GRAY)

section_label(slide, "May Daily Rates", 0.3, 4.5, 4.0)
daily_rates = [
    ("Google:", "~$1,000/day", "(down from $3,200/day at peak)", RED),
    ("Indeed:", "~$2,760/day", "(down from $4,582/day in April)",  AMBER),
    ("Reddit:", "~$26/day",    "(nearly paused — 3 evergreen campaigns)", RED),
]
for j, (ch, rate, note, col) in enumerate(daily_rates):
    add_text(slide, ch, 0.3, 4.85 + j * 0.38, 1.2, 0.32,
             font_size=9.5, bold=True, color=DARK_BLUE)
    add_text(slide, rate, 1.55, 4.85 + j * 0.38, 1.8, 0.32,
             font_size=11, bold=True, color=col)
    add_text(slide, note, 3.4, 4.88 + j * 0.38, 5.0, 0.28,
             font_size=8.5, italic=True, color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — RECOMMENDATIONS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Recommendations — Priority Action Plan")

tiers = [
    ("IMMEDIATE — This Week", RED, [
        "1. Pause Nashville Hospitality Search + Chicago Hospitality Search + "
        "Competitor Search (Google) → save ~$2,166/month at CPA $11–13",
        "2. Pause McCarran NV + Culinaire Dallas (Reddit) → save ~$4,009/month, 0 results ever",
        "3. Investigate Reno Google ($4.77 CPA) — audit keywords and audience exclusions",
    ]),
    ("REALLOCATE — Budget Shift", AMBER, [
        "4. Move Google Search budget from struggling campaigns → App Campaigns ($0.90 vs $2.42 CPA)",
        "5. Increase Google App budget in Atlanta ($0.41 CPA), Charlotte ($0.41), Orlando ($1.33)",
        "6. Reactivate Reddit with proven formats: Dishwasher Nashville ($2.03 CPR), "
        "Levy Concession Dallas ($2.32 CPR)",
    ]),
    ("INVESTIGATE — Root Cause", DARK_BLUE, [
        "7. Indeed apply completion rate — 2% anomalously low; audit ATS redirect flow",
        "8. Dallas Indeed: $11,335 spend, 65 applies — check job listing quality and landing page",
        "9. Washington DC Indeed: $3,539 with 0 platform applies — review supply availability",
    ]),
]

ty = 0.82
for tier_label, color, items in tiers:
    add_rect(slide, 0.3, ty, 12.7, 0.32, fill=color)
    add_text(slide, tier_label, 0.45, ty + 0.04, 12.4, 0.24,
             font_size=10, bold=True, color=WHITE)
    ty += 0.32
    for item in items:
        add_rect(slide, 0.3, ty, 12.7, 0.48, fill=WHITE,
                 line_color=MID_GRAY, line_width=0.3)
        add_rect(slide, 0.3, ty, 0.08, 0.48, fill=color)
        add_text(slide, item, 0.48, ty + 0.06, 12.3, 0.38,
                 font_size=8.5, color=DARK_GRAY, wrap=True)
        ty += 0.48
    ty += 0.12


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — INDEED ADS: COST PER RSVP BY LOCATION
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
slide_bg(slide)
header_bar(slide, "Indeed Ads — Cost Per RSVP by Market  |  Spend vs RSVP Attainment")

add_text(slide,
    "Cost Per RSVP = Indeed Ads Spend (May 1–19) ÷ Total RSVPs. "
    "RSVPs from RSVP Tracker 2026, snapshot May 12.",
    0.3, 0.75, 12.7, 0.25, font_size=8, italic=True, color=DARK_GRAY)

section_label(slide, "State-Level: Indeed Spend vs RSVP Attainment", 0.3, 1.02, 8.0)
headers = ["State", "Spend", "RSVPs", "Target", "Attain%", "Apply Starts", "$/RSVP", "Signal"]
col_w   = [1.9, 1.1, 1.3, 1.2, 1.0, 1.6, 1.0, 3.2]
table_y = 1.32
table_header_row(slide, headers, 0.3, table_y, col_w)

state_data = [
    ("Texas (TX)",       "$22,070", "156,034", "182,789", "85%",  "7,685", "$0.14", "✅ Efficient"),
    ("Ohio (OH)",        "$8,470",  "61,546",  "58,371",  "105%", "3,087", "$0.14", "✅ Over target"),
    ("New Jersey (NJ)",  "$4,911",  "3,604",   "13,080",  "28%",  "1,882", "$1.36", "⚠ Under-delivering"),
    ("Georgia (GA)",     "$4,385",  "64,172",  "30,892",  "208%", "1,692", "$0.07", "✅ Best CPR"),
    ("DC / Washington",  "$3,539",  "689",     "3,400",   "20%",  "1,658", "$5.14", "🔴 Worst CPR"),
    ("Tennessee (TN)",   "$2,826",  "45,965",  "52,524",  "88%",  "952",   "$0.06", "✅ Efficient"),
    ("Illinois (IL)",    "$2,047",  "12,547",  "35,297",  "36%",  "859",   "$0.16", "⚠ Gap vs target"),
    ("Nevada (NV)",      "$1,506",  "9,110",   "29,320",  "31%",  "694",   "$0.17", "⚠ Gap vs target"),
    ("Florida (FL)",     "$1,470",  "4,631",   "12,100",  "38%",  "557",   "$0.32", "⚠ Gap vs target"),
    ("Kentucky (KY)",    "$474",    "2,824",   "7,620",   "37%",  "117",   "$0.17", "⚠ Gap vs target"),
    ("N. Carolina (NC)", "$157",    "2,677",   "6,100",   "44%",  "84",    "$0.06", "✅ Efficient"),
]
cpr_col_map = {
    "$0.14": GREEN, "$1.36": AMBER, "$0.07": GREEN, "$5.14": RED,
    "$0.06": GREEN, "$0.16": None,  "$0.17": None,  "$0.32": None,
}
for i, row in enumerate(state_data):
    cpr_c  = cpr_col_map.get(row[6])
    sig_c  = GREEN if "✅" in row[7] else (RED if "🔴" in row[7] else AMBER)
    cc = [None, None, None, None, None, None, cpr_c, sig_c]
    table_data_row(slide, row, 0.3, table_y + 0.28 + i * 0.24, col_w,
                   even=(i % 2 == 0), cell_colors=cc, row_h=0.24)

callout_box(slide,
    "DC: $3,539 spent, only 689 RSVPs, 20% of target. $5.14/RSVP — 37× higher than Georgia ($0.07).",
    0.3, 4.12, 6.1, 0.45, bg=RED)
callout_box(slide,
    "Georgia ($0.07) and Tennessee ($0.06) deliver cheapest RSVPs. "
    "Ohio and Texas at or above RSVP target.",
    6.55, 4.12, 6.5, 0.45, bg=GREEN)

# Part 2 label
section_label(slide, "Client Stars & Strugglers (States with Indeed Spend >$500)", 0.3, 4.68, 12.0)

# Stars
add_rect(slide, 0.3, 4.97, 0.12, 1.85, fill=GREEN)
add_rect(slide, 0.42, 4.97, 6.1, 0.24, fill=LIGHT_BLUE, line_color=GREEN, line_width=0.5)
add_text(slide, "STARS — Exceeding RSVP Target", 0.52, 4.99, 6.0, 0.2,
         font_size=8.5, bold=True, color=DARK_BLUE)
stars = [
    ("Stord (GA)", "22,207", "458", "4,849%"),
    ("Shutterfly (TX)", "12,557", "445", "2,822%"),
    ("Asurion (TN)", "7,009", "3,053", "230%"),
    ("Foxconn (TX)", "9,076", "7,605", "119%"),
    ("Green Tokai/Yamada (OH)", "10,907", "3,876", "281%"),
]
for j, (cl, rsvp, tgt, att) in enumerate(stars):
    add_text(slide, f"{cl}: {rsvp} RSVPs / {tgt} target = {att} ✅",
             0.52, 5.22 + j * 0.3, 6.0, 0.26, font_size=8.5, color=DARK_GRAY)

# Strugglers
add_rect(slide, 6.72, 4.97, 0.12, 2.3, fill=RED)
add_rect(slide, 6.84, 4.97, 6.1, 0.24, fill=LIGHT_BLUE, line_color=RED, line_width=0.5)
add_text(slide, "STRUGGLERS — Below 50% of RSVP Target", 6.94, 4.99, 6.0, 0.2,
         font_size=8.5, bold=True, color=DARK_BLUE)
strugglers = [
    ("Indeed Flex BAU (TX)", "13,910", "29,955", "46%",  "–16,045"),
    ("OnTrac Final Mile (IL)", "5,824", "12,682", "46%",  "–6,858"),
    ("Stord (GA)", "6,992", "16,620", "42%",  "–9,628"),
    ("Stord (NV)", "4,531", "17,240", "26%",  "–12,709"),
    ("DC market (all clients)", "689", "3,400", "20%",  "–2,711"),
]
for j, (cl, rsvp, tgt, att, gap) in enumerate(strugglers):
    add_text(slide, f"{cl}: {rsvp}/{tgt} = {att}  (gap: {gap} RSVPs) 🔴",
             6.94, 5.22 + j * 0.38, 6.0, 0.34, font_size=8.5, color=DARK_GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — DATA SOURCES / CLOSING
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
add_rect(slide, 0, 0, 13.33, 7.5, fill=DARK_BLUE)
add_rect(slide, 0, 5.8, 13.33, 0.08, fill=GREEN)
add_rect(slide, 0, 5.9, 13.33, 0.05, fill=AMBER)

add_text(slide, "Data Sources & Methodology",
         1.0, 0.5, 11.0, 0.7, font_size=24, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)

sources = [
    ("Google Ads",    "Live API pull — Customer ID 7236100723 — Apr 22–May 22, 2026"),
    ("Indeed Ads",    "Dashboard export — Apr 23–May 22, 2026 + CampaignReport_Advanced May 1–21, 2026"),
    ("Reddit Ads",    "Dashboard export — Apr 22–May 21, 2026"),
    ("RSVP Data",     "RSVP Tracker 2026 — snapshot May 12, 2026"),
    ("Conversions",   "Google: app installs + form fills  |  Indeed: Apply Starts  |  Reddit: Results  |  RSVPs: RSVP Tracker"),
    ("Currency",      "All figures in USD. All periods: 2026."),
]

for j, (src, detail) in enumerate(sources):
    add_rect(slide, 1.0, 1.45 + j * 0.56, 2.1, 0.38, fill=GREEN)
    add_text(slide, src, 1.08, 1.48 + j * 0.56, 2.0, 0.3,
             font_size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, detail, 3.2, 1.48 + j * 0.56, 9.8, 0.34,
             font_size=9.5, color=WHITE)

add_text(slide, "Prepared by: Recruitment Marketing Team — Indeed Flex US",
         1.0, 5.1, 11.0, 0.35, font_size=11, color=MID_GRAY, align=PP_ALIGN.CENTER)
add_text(slide, "May 22, 2026",
         1.0, 5.45, 11.0, 0.35, font_size=11, bold=True, color=GREEN, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "/Users/claudio.santos/RM-Team-Ai/docs/reports/paid-ads-review-2026-05-22.pptx"
prs.save(out_path)
print(f"✅ Saved: {out_path}")
